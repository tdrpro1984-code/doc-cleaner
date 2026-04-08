"""Unit tests for parsers/pptx.py — PowerPoint text extraction."""
import os
import tempfile
import pytest
from pptx import Presentation
from pptx.util import Inches

from parsers.pptx import parse


@pytest.fixture
def pptx_with_text():
    """Create a PPTX file with text slides."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Body content here"
    # Second slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Second Slide"
    slide2.placeholders[1].text = "More content"
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        prs.save(f.name)
        yield f.name
    os.unlink(f.name)


@pytest.fixture
def pptx_with_table():
    """Create a PPTX file with a table."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    rows, cols = 3, 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(2))
    table = table_shape.table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Width"
    table.cell(1, 1).text = "100"
    table.cell(2, 0).text = "Height"
    table.cell(2, 1).text = "200"
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        prs.save(f.name)
        yield f.name
    os.unlink(f.name)


@pytest.fixture
def pptx_with_notes():
    """Create a PPTX file with speaker notes."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide With Notes"
    slide.placeholders[1].text = "Content"
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Remember to mention the deadline"
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        prs.save(f.name)
        yield f.name
    os.unlink(f.name)


@pytest.fixture
def pptx_empty_slide():
    """Create a PPTX with one empty slide and one text slide."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank, no text
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Only Real Slide"
    slide2.placeholders[1].text = "Content"
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        prs.save(f.name)
        yield f.name
    os.unlink(f.name)


class TestPptxTextExtraction:
    def test_slide_text(self, pptx_with_text):
        result = parse(pptx_with_text)
        assert "## Slide 1" in result
        assert "Slide Title" in result
        assert "Body content here" in result
        assert "## Slide 2" in result
        assert "Second Slide" in result

    def test_empty_slides_omitted(self, pptx_empty_slide):
        result = parse(pptx_empty_slide)
        # Empty first slide should be omitted; second slide is ## Slide 2
        assert "## Slide 2" in result
        assert "Only Real Slide" in result


class TestPptxTableExtraction:
    def test_table_as_markdown(self, pptx_with_table):
        result = parse(pptx_with_table)
        assert "Name" in result
        assert "Value" in result
        assert "|" in result  # pipe table format
        assert "---" in result  # header separator


class TestPptxSpeakerNotes:
    def test_notes_included(self, pptx_with_notes):
        result = parse(pptx_with_notes)
        assert "Notes:" in result
        assert "Remember to mention the deadline" in result

    def test_no_notes_no_label(self, pptx_with_text):
        result = parse(pptx_with_text)
        assert "Notes:" not in result


class TestPptxErrorHandling:
    def test_corrupted_file(self, tmp_path):
        bad_file = tmp_path / "corrupted.pptx"
        bad_file.write_text("not a pptx file")
        result = parse(str(bad_file))
        assert result == ""

    def test_nonexistent_file(self):
        result = parse("/nonexistent/path.pptx")
        assert result == ""
