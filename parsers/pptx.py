"""
PPTX / PPT parser — python-pptx for .pptx, textutil (macOS) fallback for .ppt.

Key features:
- Slide-level sections (## Slide N)
- Table shapes rendered as Markdown pipe tables
- Speaker notes included per slide
- Legacy .ppt via macOS textutil
"""
import os
import logging
import platform

logger = logging.getLogger(__name__)


def _table_to_markdown(table):
    """Convert a python-pptx table to a Markdown pipe table."""
    if not table.rows:
        return ""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if not rows:
        return ""
    num_cols = len(table.rows[0].cells)
    header_sep = "| " + " | ".join(["---"] * num_cols) + " |"
    rows.insert(1, header_sep)
    return "\n".join(rows)


MAX_DECOMPRESSED_SIZE = 500 * 1024 * 1024  # 500MB
MAX_SLIDES = 500


def _check_zip_size(filepath):
    """Check decompressed size of PPTX (ZIP) to prevent zip bombs."""
    import zipfile
    try:
        with zipfile.ZipFile(filepath, "r") as zf:
            total = sum(info.file_size for info in zf.infolist())
            if total > MAX_DECOMPRESSED_SIZE:
                logger.warning(
                    f"PPTX decompressed size ({total / 1024 / 1024:.0f}MB) "
                    f"exceeds limit ({MAX_DECOMPRESSED_SIZE / 1024 / 1024:.0f}MB)"
                )
                return False
    except Exception:
        pass  # not a valid zip — Presentation() will handle the error
    return True


def _parse_pptx(filepath):
    """Parse .pptx file using python-pptx."""
    from pptx import Presentation

    if not _check_zip_size(filepath):
        return ""

    try:
        prs = Presentation(filepath)
    except Exception as e:
        logger.warning(f"PPTX parse failed: {e}")
        return ""

    parts = []

    for idx, slide in enumerate(prs.slides, 1):
        if idx > MAX_SLIDES:
            logger.warning(f"PPTX slide limit reached ({MAX_SLIDES}), truncating")
            break
        slide_parts = []

        for shape in slide.shapes:
            if shape.has_table:
                md_table = _table_to_markdown(shape.table)
                if md_table:
                    slide_parts.append(md_table)
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_parts.append(text)

        # Speaker notes
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        # Skip empty slides (no text and no notes)
        if not slide_parts and not notes_text:
            continue

        section = [f"## Slide {idx}\n"]
        if slide_parts:
            section.append("\n\n".join(slide_parts))
        if notes_text:
            section.append(f"\n\n**Notes:** {notes_text}")
        parts.append("\n".join(section))

    return "\n\n".join(parts)


def _parse_ppt_textutil(filepath):
    """Parse legacy .ppt file via macOS textutil."""
    from parsers._textutil import convert_to_text
    return convert_to_text(filepath, format_label="PPT")


def parse(filepath):
    """
    Parse PowerPoint file.

    .pptx: python-pptx (slide structure, tables, notes)
    .ppt:  macOS textutil fallback (plain text only)
    """
    ext = os.path.splitext(filepath)[1].lower()

    if ext == ".pptx":
        return _parse_pptx(filepath)
    elif ext == ".ppt":
        return _parse_ppt_textutil(filepath)
    else:
        logger.warning(f"Unsupported presentation format: {ext}")
        return ""
